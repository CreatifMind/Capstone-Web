import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def match_images(pptx_path):
    prs = Presentation(pptx_path)
    print("Mapping of images to slides:")
    print("=" * 60)
    for s_idx, slide in enumerate(prs.slides):
        print(f"Slide {s_idx + 1}:")
        has_media = False
        
        # Check background fill if it has an image
        bg = slide.background
        if bg and bg.fill and bg.fill.type == 5: # Picture fill
            print("  [Background has picture/texture fill]")
            has_media = True
            
        for sh_idx, shape in enumerate(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                print(f"  Shape {sh_idx+1}: Picture (filename='{image.filename}', size={image.size}, ext='{image.ext}')")
                has_media = True
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for sub in shape.shapes:
                    if sub.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image = sub.image
                        print(f"  Shape {sh_idx+1} (Group subshape): Picture (filename='{image.filename}', size={image.size})")
                        has_media = True
            elif shape.is_placeholder:
                try:
                    if shape.has_image:
                        image = shape.image
                        print(f"  Shape {sh_idx+1}: Placeholder Picture (filename='{image.filename}', size={image.size})")
                        has_media = True
                except:
                    pass
        if not has_media:
            print("  No image shapes found.")
        print("-" * 60)

if __name__ == "__main__":
    match_images("Capstone Project.pptx")
