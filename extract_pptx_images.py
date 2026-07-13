import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_images(pptx_path, output_dir):
    os.makedirs(output_dir, exist_ok=True)
    prs = Presentation(pptx_path)
    
    img_count = 0
    for s_idx, slide in enumerate(prs.slides):
        print(f"Slide {s_idx + 1}:")
        for sh_idx, shape in enumerate(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                # Get image bytes and extension
                image_bytes = image.blob
                ext = image.ext
                filename = f"slide_{s_idx + 1}_image_{sh_idx + 1}.{ext}"
                filepath = os.path.join(output_dir, filename)
                with open(filepath, "wb") as f:
                    f.write(image_bytes)
                print(f"  Saved image: {filename}")
                img_count += 1
            elif shape.is_placeholder:
                # Some placeholders might contain images
                try:
                    if shape.has_image:
                        image = shape.image
                        image_bytes = image.blob
                        ext = image.ext
                        filename = f"slide_{s_idx + 1}_placeholder_{sh_idx + 1}.{ext}"
                        filepath = os.path.join(output_dir, filename)
                        with open(filepath, "wb") as f:
                            f.write(image_bytes)
                        print(f"  Saved placeholder image: {filename}")
                        img_count += 1
                except Exception:
                    pass
    print(f"Total images extracted: {img_count}")

if __name__ == "__main__":
    extract_images("Capstone Project.pptx", "extracted_slide_images")
